import pdfplumber
import os
from docx import Document
import pandas as pd
from pptx import Presentation
from chatbot.app.utils.text_splitter import split_text_into_chunks  # Implement this
import hashlib


def parse_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text

def parse_word(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def parse_excel(file_path):
    df = pd.read_excel(file_path)
    return df.to_string()

def parse_ppt(file_path):
    ppt = Presentation(file_path)
    text = "\n".join([slide.shapes.text for slide in ppt.slides if slide.shapes.text])
    return text

def extract_text_from_folder(folder_path, department):
    extracted_data = []

    for filename in os.listdir(folder_path):
        file_path = os.path.join(folder_path, filename)

        if filename.endswith(".pdf"):
            text = parse_pdf(file_path)
        elif filename.endswith(".docx"):
            text = parse_word(file_path)
        elif filename.endswith(".xlsx"):
            text = parse_excel(file_path)
        elif filename.endswith(".pptx"):
            text = parse_ppt(file_path)
        elif filename.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        else:
            continue  # Skip unsupported file types

        # ✅ Split long text into smaller chunks
        chunks = split_text_into_chunks(text, chunk_size=500, overlap=50)  # Example chunk size

        for chunk in chunks:
            extracted_data.append({
                "id": hashlib.md5(chunk.encode()).hexdigest(),  # Unique chunk ID
                "file_name": filename,
                "content": chunk,
                "department": department
            })

    return extracted_data
